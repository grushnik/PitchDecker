import io
import math
import base64
from datetime import date
from pathlib import Path

import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

# -----------------------------
# Assets / helpers
# -----------------------------
_ASSETS_DIR = Path(__file__).parent
EMU_PER_INCH = 914400  # pptx internal units

def _ensure_image(path_str: str) -> str:
    """Return a usable image path; if missing, create a tiny placeholder so pptx won't crash."""
    p = Path(path_str)
    if p.exists():
        return str(p)
    try:
        _png_bytes = base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR4nGNgYAAAAAMAASsJTYQAAAAASUVORK5CYII="
        )
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_bytes(_png_bytes)
    except Exception:
        pass
    return str(p)

TITLE_BG = _ensure_image(str(_ASSETS_DIR / "title_bg.png"))
BODY_BG  = _ensure_image(str(_ASSETS_DIR / "body_bg.png"))
FINAL_BG = _ensure_image(str(_ASSETS_DIR / "final_bg.png"))
MISSION_IMG = _ensure_image(str(_ASSETS_DIR / "radom_mission.png"))
TEAM_IMG    = _ensure_image(str(_ASSETS_DIR / "radom_team.png"))

def slide_size_in(prs):
    return prs.slide_width / EMU_PER_INCH, prs.slide_height / EMU_PER_INCH

def add_full_bleed_bg(slide, image_path, prs):
    slide.shapes.add_picture(image_path, Inches(0), Inches(0),
                             width=prs.slide_width, height=prs.slide_height)

# ---------- Wrapping helpers ----------
def wrap_chars(text: str, max_chars: int) -> list[str]:
    """Wrap text to lines of at most max_chars, without splitting words."""
    words = (text or "").split()
    if not words:
        return [""]
    lines, cur = [], words[0]
    for w in words[1:]:
        if len(cur) + 1 + len(w) <= max_chars:
            cur += " " + w
        else:
            lines.append(cur)
            cur = w
    lines.append(cur)
    return lines

def wrap_by_wordcount(text: str, max_words: int) -> list[str]:
    """Wrap text so each line has at most max_words (don’t split words)."""
    words = (text or "").split()
    if not words:
        return [""]
    lines, bucket = [], []
    for w in words:
        if len(bucket) < max_words:
            bucket.append(w)
        else:
            lines.append(" ".join(bucket))
            bucket = [w]
    if bucket:
        lines.append(" ".join(bucket))
    return lines

# ---- Title bar: all caps, white, centered vertically in the purple ribbon
def add_title_bar(slide, text, *, size_pt=36):
    title_text = (text or "").upper()
    tx = slide.shapes.add_textbox(Inches(0.9), Inches(0.35), Inches(11.2), Inches(1.3))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title_text
    r.font.size = Pt(size_pt)
    r.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    return tx

# ---- Bullets block: left-aligned, full-width, vertically centered then lifted 20%
def add_left_bullets_vert_center(
    slide, prs, lines, *, font_pt=28,
    left_in=0.9, right_margin_in=0.9,
    top_in=2.0, bottom_in=6.5, uplift_ratio=0.20,
    wrap_chars_limit=40
):
    raw = [l for l in (lines or ["—"]) if (l or "").strip()] or ["—"]
    wrapped_items = [wrap_chars(l, wrap_chars_limit) for l in raw]

    total_lines = sum(len(w) for w in wrapped_items)
    line_height_in = (font_pt * 1.35) / 72.0
    total_h_in = max(1, total_lines) * line_height_in

    avail_h = bottom_in - top_in
    base_top = top_in + max(0.0, (avail_h - total_h_in) / 2.0)
    shift_up = uplift_ratio * avail_h
    start_top = max(top_in, base_top - shift_up)

    slide_w_in, _ = slide_size_in(prs)
    width_in = max(1.0, slide_w_in - left_in - right_margin_in)

    tx = slide.shapes.add_textbox(Inches(left_in), Inches(start_top),
                                  Inches(width_in), Inches(total_h_in + 0.1))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    for i, chunks in enumerate(wrapped_items):
        for j, chunk in enumerate(chunks):
            p = tf.add_paragraph() if (i > 0 or j > 0) else tf.paragraphs[0]
            r = p.add_run()
            r.text = f"• {chunk}" if j == 0 else f"  {chunk}"
            r.font.size = Pt(font_pt)
            r.font.color.rgb = RGBColor(0, 0, 0)
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
    return tx

def add_center_paragraph(
    slide, prs, text, *, font_pt=30,
    box_width_in=None,  # if None, uses full width minus margins
    left_in=0.9, right_margin_in=0.9,
    top_in=2.0, bottom_in=6.5, uplift_ratio=0.20,
    wrap_chars_limit=20, align_center=True
):
    chunks = wrap_chars(text or "—", wrap_chars_limit)

    line_height_in = (font_pt * 1.35) / 72.0
    total_h_in = max(1, len(chunks)) * line_height_in
    avail_h = bottom_in - top_in
    base_top = top_in + max(0.0, (avail_h - total_h_in) / 2.0)
    shift_up = uplift_ratio * avail_h
    start_top = max(top_in, base_top - shift_up)

    slide_w_in, _ = slide_size_in(prs)
    width_in = box_width_in if box_width_in else max(1.0, slide_w_in - left_in - right_margin_in)
    # If width explicitly set, center the box horizontally
    if box_width_in:
        left_in = max(0.5, (slide_w_in - box_width_in) / 2.0)

    tx = slide.shapes.add_textbox(Inches(left_in), Inches(start_top),
                                  Inches(width_in), Inches(total_h_in + 0.2))
    tf = tx.text_frame
    tf.clear()
    for i, line in enumerate(chunks):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_pt)
        run.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER if align_center else PP_ALIGN.LEFT
    return tx

def add_footer_links_right(slide, prs, links, *, font_pt=14, right_margin_in=0.9, bottom_in=6.85):
    if not links:
        return
    slide_w_in, _ = slide_size_in(prs)
    left_in = 0.9
    width_in = max(1.0, slide_w_in - left_in - right_margin_in)
    tx = slide.shapes.add_textbox(Inches(left_in), Inches(bottom_in), Inches(width_in), Inches(0.5))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    for i, url in enumerate(links, start=1):
        if i > 1:
            sep = p.add_run(); sep.text = "  "; sep.font.size = Pt(font_pt)
        run = p.add_run()
        run.text = f"[{i}]"
        run.font.size = Pt(font_pt)
        run.font.color.rgb = RGBColor(0, 0, 128)
        try:
            run.hyperlink.address = url
        except Exception:
            pass

def target_counts(total, ratios):
    raw = {k: total * v for k, v in ratios.items()}
    floors = {k: int(math.floor(x)) for k, x in raw.items()}
    remainder = total - sum(floors.values())
    fracs = sorted([(raw[k] - floors[k], k) for k in ratios], reverse=True)
    for i in range(remainder):
        floors[fracs[i % len(fracs)][1]] += 1
    return floors

# -----------------------------
# Streamlit UI
# -----------------------------
st.set_page_config(page_title="Deck Builder inspired by 3-Minute Rule", layout="wide")
st.title("Deck Builder inspired by the 3-Minute Rule (Brant Pinvidic)")

with st.sidebar:
    st.header("Deck Header")
    project_title = st.text_input("Pitch Deck Name", "Uber for cats")
    author = st.text_input("Creator Name", "Elizaveta Grushnikova")
    place = st.text_input("Place", "Pewaukee, WI")
    date_str = st.text_input("Date", str(date.today()))
    st.markdown("---")

# Defaults
st.subheader("Hook (one or two strong lines)")
hook = st.text_area("Your hook", value="From nap to vet in one tap!", height=100)

st.subheader("But Funnel (short real-world example)")
but_funnel_default = (
    "Last month, a cat named Mittens disappeared in Milwaukee and turned up in Madison.\n"
    "Her owner opened the app, hit \"Express,\" and Uber for Cats brought her home before dinner.\n"
    "In Europe, there are about 1,000 biodigesters, largely because they receive government funding."
)
but_funnel = st.text_area("This story will appear on the 'Real-world example' slide",
                          value=but_funnel_default, height=160)

st.subheader("Enter bullets and set priority (High=5, Medium=3, None=0)")

def make_editor(title, include_link=False, key_prefix="", initial_rows=None, show_title=True, height=320, seed_rows=0):
    if show_title and title:
        st.markdown(f"**{title}**")

    cols = {
        "Bullet": pd.Series(dtype="str"),
        "High (5)": pd.Series(dtype="bool"),
        "Med (3)": pd.Series(dtype="bool"),
        "None (0)": pd.Series(dtype="bool"),
    }
    if include_link:
        cols["Link (optional)"] = pd.Series(dtype="str")
    df = pd.DataFrame(cols)

    for _ in range(max(0, seed_rows)):
        row = {"Bullet": "", "High (5)": False, "Med (3)": False, "None (0)": False}
        if include_link:
            row["Link (optional)"] = ""
        df.loc[len(df)] = row

    if initial_rows:
        for row in initial_rows:
            base = {"Bullet": row.get("Bullet", ""),
                    "High (5)": row.get("High (5)", False),
                    "Med (3)": row.get("Med (3)", False),
                    "None (0)": row.get("None (0)", False)}
            if include_link:
                base["Link (optional)"] = row.get("Link (optional)", "")
            df.loc[len(df)] = base

    df = st.data_editor(
        df,
        num_rows="dynamic",
        use_container_width=True,
        height=height,
        column_config={
            "Bullet": st.column_config.TextColumn(width="large"),
            "High (5)": st.column_config.CheckboxColumn(default=False),
            "Med (3)": st.column_config.CheckboxColumn(default=False),
            "None (0)": st.column_config.CheckboxColumn(default=False),
            "Link (optional)": st.column_config.TextColumn(width="medium") if include_link else None,
        },
        key=f"editor_{key_prefix}",
    )

    items = []
    for _, row in df.iterrows():
        text = (row.get("Bullet") or "").strip()
        if not text:
            continue
        prio = 5 if row.get("High (5)") else 3 if row.get("Med (3)") else 0
        link = (row.get("Link (optional)") or "").strip() if include_link else ""
        items.append({"text": text, "priority": prio, "link": link})
    return items

# Stacked editors (consistent spacing)
st.markdown("### WHAT IS IT?")
st.caption(
    "What do you do? What do you do well? What is it? Why is it good for? "
    "What do you want someone to do or buy? Why should they do it or buy it? What is in it for them?"
)
what_items = make_editor("", key_prefix="what", show_title=False, height=520, seed_rows=10)

st.markdown("### HOW DOES IT WORK?")
how_items = make_editor("", key_prefix="how", show_title=False, height=360, seed_rows=5)

st.markdown("### ARE YOU SURE? (include links)")
sure_items = make_editor("", include_link=True, key_prefix="sure", show_title=False, height=420, seed_rows=6)

st.markdown("### CAN YOU DO IT? (team, pilots, business case)")
cydi_items = make_editor("", key_prefix="cydi", show_title=False, height=360, seed_rows=5)

st.markdown("### DOWNSIDES (2–3 short, honest risks or trade-offs)")
downsides_items = make_editor(
    "", key_prefix="downsides", show_title=False,
    initial_rows=[{"Bullet": "Maintenance expenses spike during shedding season.", "Med (3)": True}],
    height=260, seed_rows=3
)

# WHAC weights
WHAC_RATIOS = {"WHAT": 0.50, "HOW": 0.27, "SURE": 0.15, "CYDI": 0.08}

# -----------------------------
# PPT Builder
# -----------------------------
def build_ppt(payload):
    prs = Presentation()
    blank = prs.slide_layouts[6]

    # 1) Title
    s1 = prs.slides.add_slide(blank)
    add_full_bleed_bg(s1, TITLE_BG, prs)
    title_box = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.0), Inches(1.8))
    tf = title_box.text_frame; tf.clear()
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = payload['project_title']; r.font.size = Pt(54); r.font.color.rgb = RGBColor(255,255,255)
    p.alignment = PP_ALIGN.CENTER
    meta = [payload['author'], payload['place'], payload['date']]
    box = s1.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(11.0), Inches(2.0))
    tfm = box.text_frame; tfm.clear()
    for i,line in enumerate(meta):
        par = tfm.add_paragraph() if i>0 else tfm.paragraphs[0]
        rr = par.add_run(); rr.text=line; rr.font.size=Pt(22 if i==0 else 20); rr.font.color.rgb=RGBColor(255,255,255)
        par.alignment = PP_ALIGN.CENTER

    # 2) RADOM MISSION (image centered)
    s2 = prs.slides.add_slide(blank)
    add_full_bleed_bg(s2, BODY_BG, prs)
    add_title_bar(s2, "RADOM MISSION", size_pt=36)
    slide_w_in, _ = slide_size_in(prs)
    img_max_w_in = slide_w_in - 1.8
    img_max_h_in = (6.5 - 2.0) * 0.9
    pic = s2.shapes.add_picture(MISSION_IMG, Inches(0), Inches(0))
    img_w_in = pic.width / EMU_PER_INCH
    img_h_in = pic.height / EMU_PER_INCH
    scale = min(img_max_w_in / img_w_in, img_max_h_in / img_h_in)
    pic.width = int(img_w_in * scale * EMU_PER_INCH)
    pic.height = int(img_h_in * scale * EMU_PER_INCH)
    pic.left = int((slide_w_in - (pic.width / EMU_PER_INCH)) / 2.0 * EMU_PER_INCH)
    content_top, content_bottom = 2.0, 6.5
    avail_h = content_bottom - content_top
    y_center = content_top + avail_h / 2.0 - (0.20 * avail_h)  # lifted 20%
    pic.top = int((y_center - (pic.height / EMU_PER_INCH) / 2.0) * EMU_PER_INCH)

    # helper to make bullet slides
    def make_bullets_slide(title, lines, size=28):
        s = prs.slides.add_slide(blank)
        add_full_bleed_bg(s, BODY_BG, prs)
        add_title_bar(s, title, size_pt=36)
        add_left_bullets_vert_center(
            s, prs, lines, font_pt=size,
            left_in=0.9, right_margin_in=0.9,
            top_in=2.0, bottom_in=6.5, uplift_ratio=0.20,
            wrap_chars_limit=40
        )
        return s

    # 3) WHAT IS IT?  (Top 3 WHAT)
    make_bullets_slide("WHAT IS IT?", payload["slides"]["what_top3"], size=32)

    # 4) Hook (centered; 20-char wrapping)
    s4 = prs.slides.add_slide(blank)
    add_full_bleed_bg(s4, BODY_BG, prs)
    add_title_bar(s4, "WHAT ARE WE ABOUT?", size_pt=36)
    add_center_paragraph(
        s4, prs, payload.get("hook","—"), font_pt=32,
        box_width_in=9.5,  # center the box horizontally
        top_in=2.0, bottom_in=6.5, uplift_ratio=0.20,
        wrap_chars_limit=20, align_center=True
    )

    # 5) Our team (image left, text right). Text: **max two words per line**
    s5 = prs.slides.add_slide(blank)
    add_full_bleed_bg(s5, BODY_BG, prs)
    add_title_bar(s5, "RADOM TEAM", size_pt=36)

    max_h = 6.5 - 2.0
    pic2 = s5.shapes.add_picture(TEAM_IMG, Inches(0), Inches(0))
    img2_w_in = pic2.width / EMU_PER_INCH
    img2_h_in = pic2.height / EMU_PER_INCH
    scale2 = min(5.2 / img2_w_in, max_h / img2_h_in)  # ~5.2" wide left panel
    pic2.width = int(img2_w_in * scale2 * EMU_PER_INCH)
    pic2.height = int(img2_h_in * scale2 * EMU_PER_INCH)
    pic2.left = Inches(0.9)
    avail_h2 = 6.5 - 2.0
    y_center2 = 2.0 + avail_h2 / 2.0 - (0.20 * avail_h2)
    pic2.top = int((y_center2 - (pic2.height / EMU_PER_INCH) / 2.0) * EMU_PER_INCH)

    slide_w_in, _ = slide_size_in(prs)
    gap_in = 0.6
    right_text_left = (pic2.left / EMU_PER_INCH) + (pic2.width / EMU_PER_INCH) + gap_in
    right_text_width = max(1.0, slide_w_in - right_text_left - 0.9)

    team_text = ("Diverse, resourceful, motivated team, "
                 "battle-hardened by 31 years of combined entrepreneurial experience.")
    team_lines = wrap_by_wordcount(team_text, max_words=2)  # <= 2 words per line

    # Build a vertically centered block on the right
    line_h_in = (30 * 1.35) / 72.0
    total_h_in = max(1, len(team_lines)) * line_h_in
    base_top = 2.0 + max(0.0, (avail_h2 - total_h_in) / 2.0)
    top_right = base_top - 0.20 * avail_h2

    tb = s5.shapes.add_textbox(Inches(right_text_left), Inches(top_right),
                               Inches(right_text_width), Inches(total_h_in + 0.2))
    tf = tb.text_frame; tf.clear()
    for i, line in enumerate(team_lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        r = p.add_run(); r.text = line
        r.font.size = Pt(30); r.font.color.rgb = RGBColor(0,0,0)
        p.alignment = PP_ALIGN.LEFT

    # 6) How does it work?
    make_bullets_slide("HOW DOES IT WORK?", payload["slides"]["how"], size=28)

    # 7) Downsides
    make_bullets_slide("DOWNSIDES", payload["slides"]["downsides"], size=28)

    # 8) Are you sure? (with footer links)
    s8 = prs.slides.add_slide(blank)
    add_full_bleed_bg(s8, BODY_BG, prs)
    add_title_bar(s8, "ARE YOU SURE?", size_pt=36)
    add_left_bullets_vert_center(
        s8, prs, payload["slides"]["sure_texts"], font_pt=26,
        left_in=0.9, right_margin_in=0.9, top_in=2.0, bottom_in=6.5,
        uplift_ratio=0.20, wrap_chars_limit=40
    )
    add_footer_links_right(s8, prs, payload["slides"]["sure_links"],
                           font_pt=14, right_margin_in=0.9, bottom_in=6.85)

    # 9) Can you do it?
    make_bullets_slide("CAN YOU DO IT?", payload["slides"]["cydi"], size=28)

    # 10) Real-world example (centered horizontally & vertically)
    s10 = prs.slides.add_slide(blank)
    add_full_bleed_bg(s10, BODY_BG, prs)
    add_title_bar(s10, "REAL-WORLD EXAMPLE", size_pt=36)
    add_center_paragraph(
        s10, prs, payload.get("but_funnel","—"), font_pt=28,
        box_width_in=9.5,  # center the box; slightly narrower for margins
        top_in=2.0, bottom_in=6.5, uplift_ratio=0.20,
        wrap_chars_limit=28, align_center=True
    )

    # 11) Thank you
    s11 = prs.slides.add_slide(blank)
    add_full_bleed_bg(s11, FINAL_BG, prs)
    tx2 = s11.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.0), Inches(1.5))
    tf2 = tx2.text_frame; tf2.clear()
    p2 = tf2.paragraphs[0]; r2=p2.add_run()
    r2.text="THANK YOU"; r2.font.size=Pt(44); r2.font.color.rgb=RGBColor(255,255,255); p2.alignment=PP_ALIGN.CENTER

    return prs

# -----------------------------
# Build button (drop None -> keep 50% -> WHAC targets -> High then Med)
# -----------------------------
if st.button("Build Deck"):
    # 1) Collect bullets
    all_bullets = []
    for items, cat, linkflag in [
        (what_items, "WHAT", False),
        (how_items,  "HOW",  False),
        (sure_items, "SURE", True),
        (cydi_items, "CYDI", False),
    ]:
        for idx, b in enumerate(items):
            text = (b["text"] or "").strip()
            prio = int(b["priority"] or 0)
            link = (b.get("link","").strip() if linkflag else "")
            all_bullets.append({
                "text": text,
                "priority": prio,
                "link": link,
                "category": cat,
                "idx": idx,
            })

    # 2) Exclude None(0)
    pool = [x for x in all_bullets if x["text"] and x["priority"] > 0]

    # 3) Total to keep (50% of pool)
    keep_fraction = 0.50
    kept_total = max(1, int(round(len(pool) * keep_fraction)))

    # 4) Category targets
    targets = target_counts(kept_total, {"WHAT": 0.50, "HOW": 0.27, "SURE": 0.15, "CYDI": 0.08})

    # 5) Split by category and priority tier
    by_cat = {c: {"H": [], "M": []} for c in ["WHAT","HOW","SURE","CYDI"]}
    for b in pool:
        tier = "H" if b["priority"] >= 5 else "M"
        by_cat[b["category"]][tier].append(b)

    for c in by_cat:
        by_cat[c]["H"].sort(key=lambda x: x["idx"])
        by_cat[c]["M"].sort(key=lambda x: x["idx"])

    selected = {c: [] for c in ["WHAT","HOW","SURE","CYDI"]}
    for c, cap in targets.items():
        take = []
        highs = by_cat[c]["H"][:cap]
        take.extend(highs)
        if len(take) < cap:
            need = cap - len(take)
            meds = by_cat[c]["M"][:need]
            take.extend(meds)
        selected[c] = take

    downsides = [x["text"] for x in downsides_items if (x["text"] or "").strip()][:3]
    if not downsides:
        downsides = ["The cost of plasma-based fertilizers compared to synthetic fertilizers."]

    sure_texts = [b["text"] for b in selected["SURE"]]
    sure_links = []
    for b in selected["SURE"]:
        if b.get("link"):
            url = b["link"].strip()
            if url and url not in sure_links:
                sure_links.append(url)

    what_list = [b["text"] for b in selected["WHAT"]]

    payload = {
        "project_title": project_title,
        "author": author,
        "place": place,
        "date": date_str,
        "hook": hook,
        "but_funnel": but_funnel,
        "slides": {
            "what_top3": what_list[:3],
            "what_rest": what_list[3:],
            "how":  [b["text"] for b in selected["HOW"]],
            "downsides": downsides,
            "sure_texts": sure_texts,
            "sure_links": sure_links,
            "cydi": [b["text"] for b in selected["CYDI"]],
        }
    }

    prs = build_ppt(payload)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    st.download_button(
        "Download PPTX",
        data=buf.getvalue(),
        file_name="pitch_deck_radom_style.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    st.success("Deck built! Download is ready.")
